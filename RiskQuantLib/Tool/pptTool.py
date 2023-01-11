#!/usr/bin/python
#coding = utf-8

import pandas as pd
from RiskQuantLib.Tool.wordTool import replaceParagraphContent
import pptx
from pptx import Presentation
#<import>
#</import>

def replaceTableContent(shape:pptx.shapes.graphfrm.GraphicFrame, dfInput:pd.DataFrame):
    """
    This function will replace the content of table into given dataframe. If two tables have different
    shapes, then only the common part will be changed, other cells in original table will remain unchanged.
    The change is in place, but the shape of original table won't be changed after this function is called.

    Parameters
    ----------
    shape : pptx.shapes.graphfrm.GraphicFrame
        The shape object that you want to change, usually this shape is a table.
    dfInput : pd.DataFrame
        The dataframe that the original table should be changed to.

    Returns
    -------
    None
    """
    dfInput = dfInput.reset_index().T.reset_index().T
    table = shape.table
    for row in range(min(dfInput.shape[0],len(table.rows))):
        for col in range(min(dfInput.shape[1],len(table.columns))):
            for num,paragraph in enumerate(table.cell(row,col).text_frame.paragraphs):
                try:
                    value = dfInput.iloc[row,col]
                    if isinstance(value,str):
                        paragraph.text = value
                    elif isinstance(value,float):
                        paragraph.text = format(round(value,2),",")
                    else:
                        paragraph.text = str(value)
                except Exception as e:
                    print(e)

def replaceChart(slide:pptx.slide.Slide, shape:pptx.shapes.graphfrm.GraphicFrame, picturePath:str):
    """
    This function will replace the chart in given slide into the given picture. This change is in place.

    Parameters
    ----------
    slide : pptx.slide.Slide
        The slide where the shape object lies.
    shape : pptx.shapes.graphfrm.GraphicFrame
        The shape object that you want to change, usually this shape is a chart.
    picturePath : str
        The path of picture that you want to use to replace the original chart.

    Returns
    -------
    None
    """
    tmp = shape._element
    tmp.getparent().remove(tmp)
    slide.shapes.add_picture(picturePath, shape.left, shape.top, width=shape.width, height=shape.height, )

def replacePicture(slide:pptx.slide.Slide, shape:pptx.shapes.graphfrm.GraphicFrame, picturePath:str):
    """
    This function will replace the old picture in given slide into the given picture. This change is in place.
    And the new picture has the same shape with the original one.

    Parameters
    ----------
    slide : pptx.slide.Slide
        The slide where the shape object lies.
    shape : pptx.shapes.graphfrm.GraphicFrame
        The shape object that you want to change, usually this shape is a picture.
    picturePath : str
        The path of new picture that you want to use to replace the original picture.

    Returns
    -------
    None
    """
    tmp = shape._pic
    tmp.getparent().remove(tmp)
    slide.shapes.add_picture(picturePath, shape.left, shape.top, width=shape.width, height=shape.height, )

def formatTextFrameWithTemplate(shape:pptx.shapes.graphfrm.GraphicFrame, paraDict:dict):
    """
    This function will use the content of paraDict to fill the mark in shape object. If it
    find any mark like {symbolName}, it will replace it into the value declared by paraDict.

    Parameters
    ----------
    shape : pptx.shapes.graphfrm.GraphicFrame
        The shape object that you want to change, usually this shape is a text frame.
    paraDict : dict
        The key->value pairs that declare what the symbols should be replaced to.
        The key of this dict should be string type and like '{symbolName}', and
        the value of this dict should be string type.

    Returns
    -------
    None
    """
    textFrame = shape.text_frame
    for paragraph in textFrame.paragraphs:
        replaceParagraphContent(paragraph,paraDict)

def formatTableWithTemplate(shape:pptx.shapes.graphfrm.GraphicFrame, paraDict:dict):
    """
    This function will use the content of paraDict to fill the mark in shape object. If it
    find any mark like {symbolName}, it will replace it into the value declared by paraDict.

    Parameters
    ----------
    shape : pptx.shapes.graphfrm.GraphicFrame
        The shape object that you want to change, usually this shape is a table.
    paraDict : dict
        The key->value pairs that declare what the symbols should be replaced to.
        The key of this dict should be string type and like '{symbolName}', and
        the value of this dict should be string type.

    Returns
    -------
    None
    """
    table = shape.table
    for row in range(len(table.rows)):
        for col in range(len(table.columns)):
            for num,paragraph in enumerate(table.cell(row,col).text_frame.paragraphs):
                replaceParagraphContent(paragraph,paraDict)

def formatSlide(slide:pptx.slide.Slide, tableDict:dict={},textDict:dict={},graphDict:dict={}):
    """
    This function will use the content of tableDict and textDict and graphDict to change given
    slide. The first table in the slide will be changed into tableDict[0] and the second table
    in the slide will be changed into tableDict[1], etc... The first picture in the slide will
    be changed into graphDict[0] and the second table in the slide will be changed into graphDict[1],
    etc... Any symbols like {symbolName} will be changed into given string specified by textDict.

    Parameters
    ----------
    slide : pptx.slide.Slide
        The slide which you want to change.
    tableDict :dict
        The key->value pairs that declare what the table should be replaced to. The key
        of this dict should be int number and start from 0, like 0, 1, 2... and the value
        of this dict should be pandas.DataFrame object.
    textDict : dict
        The key->value pairs that declare what the symbols should be replaced to.
        The key of this dict should be string type and like '{symbolName}', and
        the value of this dict should be string type.
    graphDict :dict
        The key->value pairs that declare what the graph should be replaced to. The key
        of this dict should be int number and start from 0, like 0, 1, 2... and the value
        of this dict should be path string of pictures.

    Returns
    -------
    None
    """
    shapes = [i for i in slide.shapes]
    tbId = 0
    imageId = 0
    for shape in shapes:
        if shape.has_table:
            if len(tableDict)!=0:
                replaceTableContent(shape, tableDict[tbId])
                tbId += 1
            if len(textDict)!=0:
                formatTableWithTemplate(shape, textDict)
        elif shape.has_text_frame and len(textDict)!=0:
            formatTextFrameWithTemplate(shape, textDict)
        elif shape.has_chart and len(graphDict)!=0:
            replaceChart(slide,shape,graphDict[imageId])
            imageId += 1
        elif type(shape) == pptx.shapes.picture.Picture and len(graphDict)!=0:
            replacePicture(slide,shape,graphDict[0])
        else:
            pass

def formatPpt(filePath:str,targetPath:str,tableDict:dict={},textDict:dict={},graphDict:dict={},slideIndex = None):
    """
    This function will use the content of tableDict and textDict and graphDict to change given
    PPT. It will iterate across every slide and modify it. The first table in the slide will be
    changed into tableDict[0] and the second table in the slide will be changed into tableDict[1],
    etc... The first picture in the slide will be changed into graphDict[0] and the second table
    in the slide will be changed into graphDict[1], etc... Any symbols like {symbolName} will be
    changed into given string specified by textDict.

    After change is done, this PPT will be saved into target file. You can also specify the slideIndex
    so that only the specified slide will be changed.

    Parameters
    ----------
    filePath : str
        The path of PPT file which you want to change.
    targetPath : str
        The path where you want to change the modified file.
    tableDict :dict
        The key->value pairs that declare what the table should be replaced to. The key
        of this dict should be int number and start from 0, like 0, 1, 2... and the value
        of this dict should be pandas.DataFrame object.
    textDict : dict
        The key->value pairs that declare what the symbols should be replaced to.
        The key of this dict should be string type and like '{symbolName}', and
        the value of this dict should be string type.
    graphDict :dict
        The key->value pairs that declare what the graph should be replaced to. The key
        of this dict should be int number and start from 0, like 0, 1, 2... and the value
        of this dict should be path string of pictures.
    slideIndex : int
        The int number i which specify the i-th slide should be changed. If you specify this
        parameter, then all the rest will not be changed.

    Returns
    -------
    None
    """
    prs = Presentation(filePath)
    if type(slideIndex)==int:
        formatSlide(prs.slides[slideIndex], tableDict, textDict, graphDict)
    else:
        [formatSlide(i,tableDict,textDict,graphDict) for i in prs.slides]
    prs.save(targetPath)

#<pptTool>
#</pptTool>